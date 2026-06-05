#!/usr/bin/env python3
"""Schule Neckertal – Signage: manifest.json aus PPTX + Gruppen-Config erzeugen.

- Liest die Gruppen-Config (Zeitplan, Overlay, Standarddauer).
- Liest die PPTX und ermittelt die SICHTBAREN Folien in Reihenfolge
  (ausgeblendete Folien `show="0"` werden ignoriert – LibreOffice rendert sie
  ebenfalls nicht, dadurch stimmt das Mapping zu den gerenderten PNGs).
- Parst pro sichtbarer Folie die Notizen auf `dauer:` / `duration:`.
- Erzeugt manifest.json passend zum Player-Schema.

Ungültige Dauerangaben brechen NICHT ab: es gilt die Standarddauer, und es wird
eine Warnung ausgegeben (Konzept Kap. 10/21.5).
"""
import argparse
import json
import os
import re
import sys
from datetime import datetime, timezone

from pptx import Presentation

DURATION_RE = re.compile(r"(?:dauer|duration)\s*:\s*([^\s]+)", re.IGNORECASE)


def parse_duration(notes, default, warnings, slide_no):
    """Dauer aus Notizen bestimmen. Gültig (>0 Ganzzahl) überschreibt Default."""
    if not notes:
        return default
    m = DURATION_RE.search(notes)
    if not m:
        return default
    raw = m.group(1).strip()
    try:
        val = int(raw)
    except ValueError:
        warnings.append(f"Folie {slide_no}: ungültige Dauer '{raw}' – Standarddauer {default}s verwendet.")
        return default
    if val <= 0:
        warnings.append(f"Folie {slide_no}: Dauer '{raw}' <= 0 – Standarddauer {default}s verwendet.")
        return default
    return val


def visible_slides(prs):
    """(Quell-Foliennummer, Notiztext) für alle sichtbaren Folien in Reihenfolge."""
    out = []
    for idx, slide in enumerate(prs.slides, start=1):
        if slide.element.get("show") == "0":      # ausgeblendet
            continue
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text or ""
        out.append((idx, notes))
    return out


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--config", required=True)
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--slides-dir", required=True, help="Verzeichnis mit slide-XXX.png")
    ap.add_argument("--output", required=True)
    ap.add_argument("--version", default=None, help="Manifest-Version (Default: UTC-Zeitstempel)")
    ap.add_argument("--slides-rel", default="slides",
                    help="Relativer Ordner der Folien im Manifest (z. B. 'slides/<version>' "
                         "für CDN-sichere, versionierte URLs)")
    ap.add_argument("--source-hash", default="",
                    help="Hash der gerenderten Folien (für Änderungserkennung im Workflow)")
    args = ap.parse_args()

    with open(args.config, encoding="utf-8") as f:
        config = json.load(f)

    default = int(config.get("defaultSlideDurationSeconds", 12))
    version = args.version or datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    warnings = []

    prs = Presentation(args.pptx)
    vis = visible_slides(prs)

    # Das PDF (Google-Export) enthält ALLE Folien in Reihenfolge -> die PNGs
    # slide-001..N entsprechen 1:1 den PPTX-Folien. Wir referenzieren je
    # sichtbarer Folie genau das PNG mit ihrer Quell-Foliennummer (übersprungene
    # Folien werden so korrekt ausgelassen).
    slides = []
    for src_no, notes in vis:
        png_name = f"slide-{src_no:03d}.png"
        if not os.path.isfile(os.path.join(args.slides_dir, png_name)):
            warnings.append(f"Folie {src_no}: gerendertes Bild fehlt ({png_name}) – übersprungen.")
            continue
        dur = parse_duration(notes, default, warnings, src_no)
        slides.append({
            "file": f"{args.slides_rel}/{png_name}",
            "durationSeconds": dur,
            "sourceSlideNumber": src_no,
        })

    manifest = {
        "groupId": config["groupId"],
        "version": version,
        "sourceHash": args.source_hash,
        "defaultSlideDurationSeconds": default,
        "schedule": config.get("schedule", {}),
        "baseLayer": {"type": "renderedPowerPoint", "slides": slides},
        "overlayLayer": config.get("overlayLayer", {}),
        "tickerLayer": config.get("tickerLayer", {"active": False, "text": ""}),
    }

    os.makedirs(os.path.dirname(os.path.abspath(args.output)), exist_ok=True)
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)

    # Statusausgabe (erscheint im Action-Log)
    print(f"manifest.json: {len(slides)} Folien, Version {version}, Standarddauer {default}s")
    for w in warnings:
        print(f"WARNUNG: {w}", file=sys.stderr)
    # Warnungen zusaetzlich als Datei (fuer spaetere Render-Statusanzeige)
    with open(os.path.join(os.path.dirname(os.path.abspath(args.output)), "render-warnings.json"), "w", encoding="utf-8") as f:
        json.dump({"version": version, "warnings": warnings}, f, indent=2, ensure_ascii=False)

    print(f"{len(warnings)} Warnung(en).")


if __name__ == "__main__":
    main()
